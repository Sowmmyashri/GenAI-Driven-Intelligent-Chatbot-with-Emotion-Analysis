import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from spellchecker import SpellChecker
from afinn import Afinn
from sklearn.metrics import accuracy_score
from transformers import pipeline
import random
from docx import Document
from pptx import Presentation


# Ensure transformers is imported correctly
try:
    from transformers import pipeline
except ImportError:
    st.error("Failed to import the transformers library. Please ensure it is installed correctly.")
    st.stop()

load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))#make sure to create .env file

# Initialize emotion analysis pipeline
emotion_analyzer = pipeline("text-classification", model="j-hartmann/emotion-english-distilroberta-base", top_k=None)
spell_checker = SpellChecker()
afinn = Afinn()

def correct_spelling(text):
    corrected_words = [spell_checker.correction(word) or word for word in text.split()]
    return ' '.join(corrected_words).strip()

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_docx_text(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        for para in doc.paragraphs:
            text += para.text
    return text

def get_txt_text(txt_docs):
    text = ""
    for txt in txt_docs:
        text += txt.read().decode('utf-8')
    return text

def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        prs = Presentation(pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_text_from_files(files):
    text = ""
    for file in files:
        if file.name.endswith(".pdf"):
            text += get_pdf_text([file])
        elif file.name.endswith(".docx"):
            text += get_docx_text([file])
        elif file.name.endswith(".txt"):
            text += get_txt_text([file])
        elif file.name.endswith(".pptx"):
            text += get_pptx_text([file])
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    return text_splitter.split_text(text)

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain(temperature=0.3):
    prompt_template = """
    You are a chatbot, and this is what you do:
    If the question is a greeting (any form of greeting in any language), respond appropriately.
    If the question is regarding your functionality or a similar question, respond appropriately based on your capabilities.
    Otherwise, follow the steps below:

    1. Read and understand the question first.
    2. Do not match the keywords in the question and the prompt. Instead, understand the context and question, and answer the question. Do not leave any small details of the context too.
    3. Focus solely on delivering informative and relevant answers. If the question is not related to the context, say "Answer is not available in the context".

    Answer the question as detailed as possible from the provided context, making sure to provide all the details. If the answer is not in the provided context, just say, "Answer is not available in the context", don't provide the wrong answer.

    Context:\n {context}\n
    Chat History:\n{chat_history}\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-1.5-flash", temperature=temperature)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "chat_history", "question"])
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

def perform_combined_sentiment_analysis(text):
    emotions = emotion_analyzer(text)
    top_emotion = max(emotions[0], key=lambda x: x['score'])
    afinn_score = afinn.score(text)
    return top_emotion['label'], top_emotion['score'], afinn_score

def get_mood(emotion_label, afinn_score):
    mood_map = {
        'joy': 'Happy 😀',
        'sadness': 'Sad 😢',
        'anger': 'Angry 😠',
        'fear': 'Concerned 😟',
        'surprise': 'Surprised 😮',
        'disgust': 'Disgusted 🤢'
    }
    mood = mood_map.get(emotion_label, 'Mixed feelings 😕')
    if afinn_score > 0:
        mood += " (Positive)"
    elif afinn_score < 0:
        mood += " (Negative)"
    return mood

def generate_response(user_question, feedback_store, chat_history):
    corrected_question = correct_spelling(user_question)
    user_emotion, user_score, afinn_score = perform_combined_sentiment_analysis(corrected_question)
    mood = get_mood(user_emotion, afinn_score)

    chat_history.append({
        "role": "user",
        "content": f"{user_question}\n\n_Emotion: {user_emotion} (score: {user_score:.2f})\n_AFINN Score: {afinn_score:.2f}\n_Mood: {mood}_"
    })

    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question, "chat_history": chat_history[-5:]}, return_only_outputs=True)

    response_text = response["output_text"]
    chat_history.append({"role": "bot", "content": response_text})

    feedback_store[corrected_question] = {
        "responses": [response_text],
        "feedback": None
    }

    return chat_history

def update_feedback(feedback_store, user_question, feedback):
    if user_question in feedback_store:
        feedback_store[user_question]["feedback"] = feedback

def regenerate_response(user_question, feedback_store, chat_history):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain(temperature=random.uniform(0.7, 1.0))  # Increase temperature for more variability
    response = chain({"input_documents": docs, "question": user_question, "chat_history": chat_history[-5:]}, return_only_outputs=True)
    new_response_text = response["output_text"]
    
    if user_question in feedback_store:
        feedback_store[user_question]["responses"].append(new_response_text)
    else:
        feedback_store[user_question] = {
            "responses": [new_response_text],
            "feedback": None
        }

    # Remove the last bot response if it exists before adding the new one
    if chat_history and chat_history[-1]["role"] == "bot":
        chat_history.pop()

    chat_history.append({"role": "bot", "content": new_response_text})
    return chat_history

def evaluate_sentiment_analysis(model, test_data, true_labels):
    predicted_labels = []
    for text in test_data:
        emotions = model(text)
        top_emotion = max(emotions[0], key=lambda x: x['score'])
        predicted_labels.append(top_emotion['label'])

    accuracy = accuracy_score(true_labels, predicted_labels) if predicted_labels else 0.0
    return accuracy

def main():
    st.set_page_config(page_title="ChatBot")

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "feedback_store" not in st.session_state:
        st.session_state.feedback_store = {}
    if "edit_index" not in st.session_state:
        st.session_state.edit_index = -1

    def edit_message(index):
        st.session_state.edit_index = index
        st.experimental_rerun()

    def save_edited_message(index, new_message):
        if index >= 0 and index < len(st.session_state.chat_history):
            user_message = new_message
            bot_response_index = index + 1 if index + 1 < len(st.session_state.chat_history) else None
            if bot_response_index and st.session_state.chat_history[bot_response_index]["role"] == "bot":
                st.session_state.chat_history = st.session_state.chat_history[:index]
                st.session_state.chat_history = generate_response(user_message, st.session_state.feedback_store, st.session_state.chat_history)
            st.session_state.edit_index = -1
            st.experimental_rerun()

    for i, chat in enumerate(st.session_state.chat_history):
        with st.chat_message(chat["role"]):
            if chat["role"] == "user" and st.session_state.edit_index == i:
                new_message = st.text_area("Edit your message", chat["content"])
                if st.button("Save", key=f"save_{i}"):
                    save_edited_message(i, new_message)
            else:
                st.markdown(chat["content"])
                if chat["role"] == "user":
                    if st.button("Edit", key=f"edit_{i}"):
                        edit_message(i)

    if prompt := st.chat_input("Type your message here..."):
        # Check for follow-up questions
        last_user_message = st.session_state.chat_history[-2]["content"] if len(st.session_state.chat_history) > 1 else None
        last_bot_response = st.session_state.chat_history[-1]["content"] if len(st.session_state.chat_history) > 0 else None

        if last_user_message and last_bot_response:
            # Determine if the current prompt is a follow-up question
            follow_up_keywords = ["more information", "what about", "details on", "clarify", "expand on"]
            is_follow_up = any(keyword in prompt.lower() for keyword in follow_up_keywords)

            if is_follow_up:
                # Use the previous context and question to generate a new response
                st.session_state.chat_history = regenerate_response(last_user_message, st.session_state.feedback_store, st.session_state.chat_history)
                st.experimental_rerun()
            else:
                # Process as a new question and generate a response
                st.session_state.chat_history = generate_response(prompt, st.session_state.feedback_store, st.session_state.chat_history)
                st.experimental_rerun()
        else:
            # No previous context, treat as a new question
            st.session_state.chat_history = generate_response(prompt, st.session_state.feedback_store, st.session_state.chat_history)
            st.experimental_rerun()

    if st.session_state.chat_history:
        last_user_message = st.session_state.chat_history[-2]["content"] if len(st.session_state.chat_history) > 1 else None
        last_bot_response = st.session_state.chat_history[-1]["content"] if len(st.session_state.chat_history) > 0 else None

        st.write("Was this response helpful?")
        col1, col2, col3 = st.columns(3)
        with col1:
            if st.button("👍"):
                update_feedback(st.session_state.feedback_store, last_user_message, "positive")
        with col2:
            if st.button("👎"):
                update_feedback(st.session_state.feedback_store, last_user_message, "negative")
        with col3:
            if st.button("Regenerate"):
                if last_user_message:
                    st.session_state.chat_history = regenerate_response(last_user_message, st.session_state.feedback_store, st.session_state.chat_history)
                    st.experimental_rerun()

    with st.sidebar:
        st.title("Menu:")
        docs = st.file_uploader("Upload your Files and Click on the Submit & Process Button", accept_multiple_files=True)
        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text = get_text_from_files(docs)
                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.success("Done")

    # Example test data and labels for accuracy evaluation
    test_data = ["I feel happy today.", "This news makes me sad.", "I am angry about the situation."]
    true_labels = ['joy', 'sadness', 'anger']
    accuracy = evaluate_sentiment_analysis(emotion_analyzer, test_data, true_labels)
    st.write(f"Sentiment Analysis Accuracy: {accuracy:.2f}")

if _name_ == "_main_":
    main()